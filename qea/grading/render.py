"""Render produced deliverable files to page images + extracted text.

LibreOffice-headless converts office files to PDF; PyMuPDF rasterizes pages to PNG.
Text is extracted per type for the text-only ablation and as judge context. If
LibreOffice or a parser is unavailable, the file degrades to text-only (logged),
never crashing the run.
"""
from __future__ import annotations

import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path

SUPPORTED = {".xlsx", ".pptx", ".docx", ".pdf"}
_MAC_APP = "/Applications/LibreOffice.app/Contents/MacOS/soffice"


@dataclass
class RenderedDeliverable:
    text: str                                  # the agent's final message
    extracted_text: str                        # text pulled from produced files
    images: list = field(default_factory=list) # list[Path] of PNG page images
    degraded: list = field(default_factory=list)  # human-readable degrade notes


def _soffice() -> str | None:
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    if Path(_MAC_APP).exists():  # brew --cask installs the .app, sometimes off PATH
        return _MAC_APP
    return None


def _to_pdf(f: Path, out_dir: Path) -> Path | None:
    if f.suffix.lower() == ".pdf":
        return f
    exe = _soffice()
    if not exe:
        return None
    out_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run([exe, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(f)],
                   check=True, capture_output=True, timeout=120)
    pdf = out_dir / (f.stem + ".pdf")
    return pdf if pdf.exists() else None


def _pdf_to_pngs(pdf: Path, out_dir: Path, limit: int) -> list:
    import fitz  # PyMuPDF
    out_dir.mkdir(parents=True, exist_ok=True)
    imgs: list = []
    doc = fitz.open(pdf)
    try:
        for i, page in enumerate(doc):
            if len(imgs) >= limit:
                break
            pix = page.get_pixmap(dpi=110)
            png = out_dir / f"{pdf.stem}_p{i + 1}.png"
            pix.save(str(png))
            imgs.append(png)
    finally:
        doc.close()
    return imgs


def _extract_text(f: Path) -> str:
    ext = f.suffix.lower()
    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(f, read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"# sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    out.append("\t".join(cells))
        return "\n".join(out)
    if ext == ".pptx":
        from pptx import Presentation
        out = []
        for i, slide in enumerate(Presentation(f).slides):
            out.append(f"# slide {i + 1}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    out.append(shape.text_frame.text)
        return "\n".join(out)
    if ext == ".docx":
        from docx import Document
        return "\n".join(p.text for p in Document(f).paragraphs)
    if ext == ".pdf":
        import fitz
        doc = fitz.open(f)
        try:
            return "\n".join(page.get_text() for page in doc)
        finally:
            doc.close()
    return ""


def render(final_text: str, files, out_dir, *, max_images: int = 8) -> RenderedDeliverable:
    out_dir = Path(out_dir)
    texts: list[str] = []
    images: list = []
    degraded: list[str] = []
    for raw in files:
        f = Path(raw)
        if f.suffix.lower() not in SUPPORTED:
            continue
        try:
            t = _extract_text(f)
            if t.strip():
                texts.append(f"=== {f.name} ===\n{t}")
        except Exception as exc:  # noqa: BLE001
            degraded.append(f"{f.name}: text-extract failed ({type(exc).__name__}: {exc})")
        if len(images) >= max_images:
            continue
        try:
            pdf = _to_pdf(f, out_dir)
            if pdf is None:
                degraded.append(f"{f.name}: LibreOffice unavailable, render skipped (text-only)")
                continue
            images.extend(_pdf_to_pngs(pdf, out_dir, max_images - len(images)))
        except Exception as exc:  # noqa: BLE001
            degraded.append(f"{f.name}: render failed ({type(exc).__name__}: {exc})")
    return RenderedDeliverable(final_text, "\n\n".join(texts), images, degraded)
